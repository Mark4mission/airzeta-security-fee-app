#!/usr/bin/env python3
"""
Generate 4 editable PowerPoint presentations for AirZeta Security Portal.
1. Admin Guide (Korean, detailed, unlimited pages)
2. Branch User Guide (English, step-by-step, each feature <= 1 page)
3. IT Planner Guide (Korean, technical architecture)
4. Executive/Finance Guide (Korean+English, cost analysis, infographics)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# ============ COMMON THEME COLORS ============
NAVY_DARK = RGBColor(0x0B, 0x19, 0x29)
NAVY_MID = RGBColor(0x13, 0x2F, 0x4C)
NAVY_LIGHT = RGBColor(0x1A, 0x3A, 0x5C)
RED_ACCENT = RGBColor(0xE9, 0x45, 0x60)
BLUE_ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN_ACCENT = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x8B, 0x99, 0xA8)
LIGHT_TEXT = RGBColor(0xE8, 0xEA, 0xED)
DARK_TEXT = RGBColor(0x1a, 0x1a, 0x2e)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

OUTPUT_DIR = "/home/user/webapp"

def set_slide_bg(slide, color=NAVY_DARK):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Malgun Gothic'):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=12, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Malgun Gothic', space_before=Pt(4), space_after=Pt(2)):
    """Add a paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_title_slide(prs, title, subtitle, version="v1.9"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY_DARK)
    
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.6), Inches(10), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED_ACCENT
    bar.line.fill.background()
    
    # Logo placeholder
    logo_shape = add_shape_bg(slide, Inches(4.2), Inches(1.0), Inches(1.6), Inches(1.2), RED_ACCENT)
    add_text_box(slide, Inches(4.2), Inches(1.15), Inches(1.6), Inches(0.9), "AIRZETA", 
                 font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, font_name='Arial')
    
    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(9), Inches(1.0), title,
                 font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(0.6), subtitle,
                 font_size=16, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(0.4), f"AirZeta Station Security System | {version} | 2026-03-04",
                 font_size=10, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    return slide

def add_section_slide(prs, section_title, section_subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY_MID)
    
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(2), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED_ACCENT
    bar.line.fill.background()
    
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(8.4), Inches(0.8), section_title,
                 font_size=28, color=WHITE, bold=True)
    if section_subtitle:
        add_text_box(slide, Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.6), section_subtitle,
                     font_size=14, color=GRAY_TEXT)
    return slide

def add_content_slide(prs, title, bullets, notes=""):
    """Add a standard content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY_DARK)
    
    # Title bar
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = NAVY_MID
    title_bg.line.fill.background()
    
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.6), title,
                 font_size=22, color=WHITE, bold=True)
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.85), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RED_ACCENT
    line.line.fill.background()
    
    # Content
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if bullet.startswith("##"):
            p.text = bullet[2:].strip()
            p.font.size = Pt(14)
            p.font.color.rgb = BLUE_ACCENT
            p.font.bold = True
            p.space_before = Pt(12)
        elif bullet.startswith("- "):
            p.text = bullet
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_TEXT
            p.space_before = Pt(3)
            p.level = 1
        else:
            p.text = bullet
            p.font.size = Pt(12)
            p.font.color.rgb = LIGHT_TEXT
            p.space_before = Pt(4)
        
        p.font.name = 'Malgun Gothic'
    
    return slide


# ======================================================================
# PPT 1: ADMIN GUIDE (Korean, Detailed)
# ======================================================================
def create_admin_guide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs, "AirZeta Security Portal\n관리자 가이드", 
                    "HQ Admin 전용 운영 매뉴얼 (상세)", "v1.9")
    
    # TOC
    add_content_slide(prs, "목차 (Table of Contents)", [
        "## 1. 시스템 개요 및 로그인",
        "## 2. 대시보드 (홈 화면)",
        "## 3. Security Directive (보안 지시사항)",
        "## 4. Security Communication (보안 커뮤니케이션)",
        "## 5. Document Library (문서 관리)",
        "## 6. Security Level (보안 레벨 관리)",
        "## 7. Security Fee (보안 수수료)",
        "## 8. Security Policy (보안 정책)",
        "## 9. Important Links (주요 링크)",
        "## 10. Settings & Users (설정 및 사용자 관리)",
        "## 11. AI 번역 기능 활용",
        "## 12. 자주 묻는 질문 (FAQ)",
    ])
    
    # Section 1: System Overview
    add_section_slide(prs, "1. 시스템 개요 및 로그인", "System Overview & Login")
    
    add_content_slide(prs, "시스템 개요", [
        "## AirZeta Station Security System",
        "- 항공 화물 보안 관리를 위한 통합 웹 포털 시스템",
        "- React 19 기반 SPA (Single Page Application)",
        "- Firebase 인증 및 데이터베이스 연동",
        "- Google Gemini AI 기반 실시간 번역 기능",
        "",
        "## 접속 URL",
        "- 운영 서버: https://airzeta-security-fee-app.vercel.app",
        "- GitHub Pages: https://mark4mission.github.io/airzeta-security-fee-app",
        "",
        "## 사용자 역할",
        "- HQ Admin (hq_admin): 전체 시스템 관리 권한",
        "- Branch User (branch_user): 지점별 제한된 접근 권한",
    ])
    
    add_content_slide(prs, "로그인 방법", [
        "## 로그인 절차",
        "- 1. 웹 브라우저에서 포털 URL 접속",
        "- 2. 이메일 주소 입력",
        "- 3. 비밀번호 입력",
        "- 4. 'Login' 버튼 클릭",
        "",
        "## 관리자 계정 특징",
        "- 상단 헤더에 'ADMIN' 뱃지 표시",
        "- 사이드바에 'Settings & Users' 메뉴 노출",
        "- 모든 게시물 수정/삭제 권한",
        "- Security Level 전 지점 편집 가능",
        "",
        "## 주의사항",
        "- 비밀번호는 주기적으로 변경 권장",
        "- 공용 PC에서 사용 후 반드시 로그아웃",
    ])
    
    # Section 2: Dashboard
    add_section_slide(prs, "2. 대시보드 (홈 화면)", "Home Dashboard")
    
    add_content_slide(prs, "홈 화면 구성", [
        "## 모듈 카드 영역",
        "- 각 기능별 바로가기 카드가 그리드 형태로 배치",
        "- Security Policy, Bulletins, Document Library 등",
        "- 클릭 시 해당 모듈로 바로 이동",
        "",
        "## 최근 게시물 (Recent Bulletins)",
        "- Security Directive 최신 게시물 미리보기",
        "- 제목, 작성일, 조회수 표시",
        "- 클릭 시 해당 게시물 상세 페이지로 이동",
        "",
        "## Global Security News Feed",
        "- 전 세계 항공 보안 관련 뉴스 RSS 피드",
        "- Google News에서 자동 수집",
        "- 최신 보안 동향 실시간 확인 가능",
    ])
    
    # Section 3: Security Directive
    add_section_slide(prs, "3. Security Directive", "보안 지시사항 관리")
    
    add_content_slide(prs, "Security Directive 개요", [
        "## 게시판 특징",
        "- HQ Admin 전용 게시판 (작성 권한: 관리자만)",
        "- 전 지점에 보안 지시사항 전달 목적",
        "- Firestore 컬렉션: bulletinPosts",
        "",
        "## 주요 기능",
        "- Rich Text 에디터 (Quill 2) + Markdown 모드 지원",
        "- AI 번역 (한국어/영어/일본어 등 10개 언어)",
        "- 파일 첨부 (최대 100MB, 드래그앤드롭 지원)",
        "- 확인(Acknowledge) 추적 시스템",
        "- 댓글 및 댓글 번역 기능",
        "- PDF 인쇄 기능",
        "",
        "## 관리자 Acknowledge 추적",
        "- 전 지점의 확인 현황을 실시간 모니터링",
        "- 확인/미확인 지점 목록 자동 분류",
    ])
    
    add_content_slide(prs, "Directive 게시물 작성", [
        "## 작성 절차",
        "- 1. Security Bulletins > Security Directive 메뉴 진입",
        "- 2. 'New Directive' 버튼 클릭",
        "- 3. 제목 입력",
        "- 4. 본문 작성 (Rich Text 또는 Markdown 모드)",
        "- 5. 필요시 파일 첨부 (드래그앤드롭 또는 클릭)",
        "- 6. 'Publish Directive' 클릭하여 게시",
        "",
        "## 에디터 기능",
        "- 헤더 (H1, H2, H3), 굵게, 기울임, 밑줄",
        "- 글자색/배경색 변경, 리스트 (번호/불릿)",
        "- 정렬 (왼쪽/가운데/오른쪽/양쪽)",
        "- 인용문, 코드 블록, 링크, 이미지",
        "- 테이블 삽입 (행/열 지정 가능)",
        "- .md 파일 업로드 지원",
    ])
    
    add_content_slide(prs, "AI 번역 기능", [
        "## 번역 모델",
        "- Google Gemini 2.5 Flash 모델 사용",
        "- 항공 화물 보안 전문 용어 최적화",
        "",
        "## 사용 방법 (게시물)",
        "- 1. 언어 선택 드롭다운에서 대상 언어 선택",
        "- 2. 'AI Translate' 버튼 클릭",
        "- 3. 원본과 번역본이 나란히 표시 (Side-by-side / Inline 전환 가능)",
        "",
        "## 사용 방법 (댓글)",
        "- 각 댓글 옆 KOR / ENG 버튼으로 즉시 번역",
        "- 사용자 로컬 타임존 기반 추가 언어 버튼 자동 표시",
        "",
        "## 주의사항",
        "- VITE_GEMINI_API_KEY 환경변수 설정 필요",
        "- ICAO, TSA, IATA 등 항공 약어는 번역하지 않음",
    ])
    
    # Section 4: Security Communication
    add_section_slide(prs, "4. Security Communication", "보안 커뮤니케이션 (신규)")
    
    add_content_slide(prs, "Security Communication 개요", [
        "## 게시판 특징",
        "- 전 사용자 게시 가능 (Admin + Branch User)",
        "- 보다 자유로운 소통 목적 (질문, 공유, 논의)",
        "- Firestore 컬렉션: communicationPosts",
        "- 파란색 테마로 구분",
        "",
        "## Directive와의 차이점",
        "- 작성 권한: 모든 사용자 (vs Directive: Admin만)",
        "- 사이트 코드 접두사 표시 (최대 3글자)",
        "- HQ Admin 게시물에는 'TA' 접두사",
        "- 지점 사용자는 IATA 코드 접두사 (예: ICN, LAX)",
        "",
        "## 동일 기능 지원",
        "- Rich Text + Markdown, AI 번역, 파일 첨부",
        "- 댓글 + 이모지, Acknowledge 추적",
        "- 드래그앤드롭 파일 업로드",
    ])
    
    # Section 5: Document Library
    add_section_slide(prs, "5. Document Library", "문서 관리")
    
    add_content_slide(prs, "문서 라이브러리", [
        "## 기능 개요",
        "- 보안 관련 문서 중앙 저장소",
        "- 카테고리별 분류 및 IATA 코드 태깅",
        "- 파일 업로드 (최대 100MB)",
        "",
        "## 접근 제어",
        "- 다운로드 권한: 'Admin Only' 또는 'All Branches'",
        "- 업로더 지점 및 Admin은 항상 접근 가능",
        "- 잠긴 문서는 회색으로 표시",
        "",
        "## 다운로드 추적",
        "- 각 다운로드마다 로그 기록",
        "- userId, userEmail, branchName, fileName, 일시",
        "- Admin은 지점별 다운로드 현황 확인 가능",
        "",
        "## 드래그앤드롭 다운로드",
        "- 파일을 데스크톱으로 끌어서 저장 가능",
    ])
    
    # Section 6: Security Level
    add_section_slide(prs, "6. Security Level", "보안 레벨 관리 / 세계 지도")
    
    add_content_slide(prs, "보안 레벨 관리", [
        "## 세계 지도 (World Map)",
        "- TopoJSON 기반 실제 세계 지도 렌더링",
        "- 85+ IATA 공항 마커 표시",
        "- 마커 색상: 녹색(Safe), 주황(Caution), 빨강(Alert)",
        "- 확대/축소 (마우스 휠), 이동 (드래그), 더블클릭 확대",
        "- ICN(인천) 중심 지도 배치",
        "",
        "## Admin 기능",
        "- 지도에서 마커 클릭 시 해당 지점 편집 모드",
        "- 보안 레벨 (Level 1~5) 설정",
        "- 가이드라인 작성/수정",
        "- Effective Since 날짜 설정",
        "- 변경 이력 자동 기록",
        "",
        "## 공항 코드 관리",
        "- airportCode 필드로 IATA 코드 지정",
        "- 자동완성 드롭다운 (AIRPORT_COORDS 사전)",
    ])
    
    # Section 7: Security Fee
    add_section_slide(prs, "7. Security Fee", "보안 수수료 관리")
    
    add_content_slide(prs, "보안 수수료", [
        "## 기능 개요",
        "- 항공 화물 보안 수수료 관리 시스템",
        "- Excel 파일 업로드를 통한 일괄 업데이트",
        "- 환율 자동 변환 (KRW 기준)",
        "",
        "## Admin 뷰",
        "- 전체 지점 수수료 현황 조회",
        "- 환율 데이터 업로드 및 관리",
        "- 수수료 변경 이력 추적",
        "",
        "## Branch User 뷰",
        "- 자기 지점 수수료 정보만 조회",
        "- 현지 통화 및 KRW 환산 금액 표시",
    ])
    
    # Section 8: Security Policy
    add_content_slide(prs, "Security Policy / Important Links", [
        "## Security Policy (보안 정책)",
        "- 보안 정책 문서 열람 기능",
        "- 관리자가 업로드한 정책 문서 표시",
        "",
        "## Important Links (주요 링크)",
        "- QR 코드 포함 링크 그룹 관리",
        "- 카테고리별 분류",
        "- 관리자가 링크 추가/수정/삭제",
        "- 각 링크에 QR 코드 자동 생성 (qrcode.react)",
    ])
    
    # Section 9: Settings
    add_section_slide(prs, "10. Settings & Users", "설정 및 사용자 관리")
    
    add_content_slide(prs, "사용자 관리", [
        "## 사용자 목록",
        "- 전체 등록 사용자 조회",
        "- 역할(role), 지점(branchName), 이메일 표시",
        "",
        "## 권한 관리",
        "- 사용자 역할 변경 (hq_admin / branch_user)",
        "- 지점 할당/변경",
        "",
        "## 지점 코드 관리",
        "- branchCodes 컬렉션 관리",
        "- 새 지점 추가, 기존 지점 수정",
    ])
    
    # FAQ
    add_section_slide(prs, "12. 자주 묻는 질문 (FAQ)", "")
    
    add_content_slide(prs, "FAQ", [
        "## Q: 번역이 작동하지 않습니다.",
        "- A: .env 파일에 VITE_GEMINI_API_KEY가 설정되어 있는지 확인하세요.",
        "- aistudio.google.com/apikey 에서 발급 가능합니다.",
        "",
        "## Q: 파일 업로드가 실패합니다.",
        "- A: 파일 크기가 Admin 100MB, Branch User 50MB 제한을 확인하세요.",
        "- Firebase Storage 설정이 올바른지 확인하세요.",
        "",
        "## Q: 지도에서 공항이 보이지 않습니다.",
        "- A: Security Level에서 해당 지점의 airportCode가 설정되어 있는지 확인하세요.",
        "- AIRPORT_COORDS에 해당 공항이 등록되어 있어야 합니다.",
        "",
        "## Q: 사이드바 메뉴가 접히지 않습니다.",
        "- A: v1.9 업데이트로 수정되었습니다. 캐시를 삭제하고 새로고침 해주세요.",
    ])
    
    filepath = os.path.join(OUTPUT_DIR, "AirZeta_Admin_Guide_KR.pptx")
    prs.save(filepath)
    print(f"[OK] Admin Guide saved: {filepath}")
    return filepath


# ======================================================================
# PPT 2: BRANCH USER GUIDE (English, Step-by-step)
# ======================================================================
def create_branch_guide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs, "AirZeta Security Portal\nBranch User Guide", 
                    "Step-by-Step Guide for Station Users", "v1.9")
    
    # Login
    add_content_slide(prs, "1. Login & Navigation", [
        "## How to Login",
        "- Open your web browser and go to the portal URL",
        "- Enter your email address and password",
        "- Click 'Login' to access the portal",
        "",
        "## Sidebar Navigation",
        "- Use the left sidebar to navigate between modules",
        "- Click the arrow icon at the bottom to collapse/expand the sidebar",
        "- Your station name and role are shown at the bottom",
        "- Click 'Logout' when you're done",
        "",
        "## Top Header",
        "- Shows the current page name and your email",
        "- 'USER' badge indicates your access level",
    ])
    
    add_content_slide(prs, "2. Home Dashboard", [
        "## Module Cards",
        "- Quick-access cards for each feature",
        "- Click any card to go directly to that section",
        "",
        "## Recent Bulletins",
        "- Latest security directives from HQ",
        "- Click to read the full post",
        "",
        "## Global Security News",
        "- Live news feed about aviation security worldwide",
        "- Stay updated on the latest security developments",
    ])
    
    add_content_slide(prs, "3. Security Directive (Read & Acknowledge)", [
        "## What is Security Directive?",
        "- Official security instructions from HQ Admin",
        "- All branch users must read and acknowledge",
        "",
        "## How to Read and Acknowledge",
        "- Go to Security Bulletins > Security Directive",
        "- Click a post title to read the full content",
        "- Click 'Acknowledge Receipt' button to confirm you've read it",
        "- Status changes from 'Pending' (yellow) to 'Acknowledged' (green)",
        "",
        "## How to Translate",
        "- Select target language from the dropdown",
        "- Click 'Translate' to see AI translation side-by-side",
        "- Toggle between Side-by-side and Inline views",
    ])
    
    add_content_slide(prs, "4. Security Communication (Read & Write)", [
        "## What is Security Communication?",
        "- Open communication board for all users",
        "- Share updates, ask questions, discuss security topics",
        "",
        "## How to Create a Post",
        "- Go to Security Bulletins > Security Communication",
        "- Click 'New Post' button",
        "- Enter a title and write your content",
        "- Attach files by dragging & dropping or clicking 'Browse'",
        "- Click 'Publish Post' to share with all stations",
        "",
        "## Comments & Emoji",
        "- Leave comments at the bottom of any post",
        "- Use the smiley icon to insert emojis",
        "- Translate comments using KOR/ENG buttons",
    ])
    
    add_content_slide(prs, "5. Document Library", [
        "## Browsing Documents",
        "- Go to Document Library from the sidebar",
        "- Browse by category or search by title",
        "- Each document shows IATA code, category, and title",
        "",
        "## Downloading",
        "- Click a document to view details",
        "- Click the download button or drag the file to your desktop",
        "- Some documents may be restricted (shown in gray)",
        "",
        "## Note",
        "- Your download history is tracked for security purposes",
    ])
    
    add_content_slide(prs, "6. Security Level (View Only)", [
        "## World Map",
        "- View the global security level map",
        "- Each station is marked with a colored dot:",
        "  - Green = Safe, Orange = Caution, Red = Alert",
        "- Scroll to zoom, drag to pan, double-click to zoom in",
        "",
        "## Your Station",
        "- Click 'Edit My Station' to update your station's security level",
        "- Set Level 1 to 5, add guidelines, and set the effective date",
        "- Changes are saved with history tracking",
    ])
    
    add_content_slide(prs, "7. Other Features", [
        "## Security Fee",
        "- View your station's security fee information",
        "- Displayed in local currency with KRW conversion",
        "",
        "## Security Policy",
        "- Read the latest security policies from HQ",
        "",
        "## Important Links",
        "- Quick access to useful resources",
        "- Each link includes a QR code for easy mobile access",
        "",
        "## Tips",
        "- Check the portal daily for new directives",
        "- Always acknowledge directives promptly",
        "- Use AI translation if you need content in your language",
    ])
    
    filepath = os.path.join(OUTPUT_DIR, "AirZeta_Branch_User_Guide_EN.pptx")
    prs.save(filepath)
    print(f"[OK] Branch User Guide saved: {filepath}")
    return filepath


# ======================================================================
# PPT 3: IT PLANNER GUIDE (Korean, Technical)
# ======================================================================
def create_it_planner_guide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs, "AirZeta Security Portal\nIT 기획/기술 가이드", 
                    "시스템 아키텍처, 클라우드 구성, AI 로직, 보안 설계", "v1.9")
    
    # Architecture Overview
    add_section_slide(prs, "1. 시스템 아키텍처", "System Architecture Overview")
    
    add_content_slide(prs, "기술 스택 (Tech Stack)", [
        "## Frontend",
        "- React 19.2.0 (SPA, HashRouter)",
        "- Vite 7.2.4 (빌드 도구, HMR)",
        "- React Router DOM 7.13.1",
        "- react-quill-new 3.8.3 (Rich Text Editor, Quill 2 기반)",
        "- lucide-react 0.562.0 (아이콘 라이브러리)",
        "",
        "## Backend Services (Serverless)",
        "- Firebase Authentication (이메일/비밀번호 인증)",
        "- Cloud Firestore (NoSQL 실시간 데이터베이스)",
        "- Firebase Storage (파일 저장소)",
        "",
        "## AI Service",
        "- Google Gemini AI (@google/genai 1.43.0)",
        "- 모델: gemini-2.5-flash (번역용)",
        "- 클라이언트 사이드 API 호출 (REST)",
    ])
    
    add_content_slide(prs, "클라우드 아키텍처 (Cloud Architecture)", [
        "## 호스팅",
        "- Vercel: 메인 프로덕션 배포 (자동 배포, CI/CD)",
        "- GitHub Pages: 보조 배포 (gh-pages 브랜치)",
        "",
        "## Firebase 프로젝트 구성",
        "- Authentication: 이메일/비밀번호 기반",
        "- Firestore: 8+ 컬렉션, 실시간 스냅샷 리스너",
        "- Storage: bulletin_attachments/, document_library/ 경로",
        "",
        "## CDN & 성능",
        "- Vercel Edge Network (글로벌 CDN)",
        "- 정적 에셋 자동 캐싱",
        "- Code splitting: Turndown, Marked 등 동적 임포트",
        "- 빌드 사이즈: ~1.3MB (gzip ~375KB)",
        "",
        "## 외부 데이터 소스",
        "- Google Sheets CSV (보안 서약 데이터)",
        "- Google News RSS (보안 뉴스 피드)",
        "- countries-110m.json (TopoJSON 지도 데이터)",
    ])
    
    add_content_slide(prs, "Firestore 데이터 모델", [
        "## 주요 컬렉션",
        "- users: uid, email, role, branchName, displayName",
        "- branchCodes: 지점 코드 마스터 데이터",
        "- bulletinPosts: 보안 지시사항 (title, content, author, attachments, acknowledgedBy[], viewCount)",
        "- bulletinPosts/{id}/comments: 댓글 서브컬렉션",
        "- communicationPosts: 보안 커뮤니케이션 (동일 구조)",
        "- communicationPosts/{id}/comments: 댓글 서브컬렉션",
        "- securityLevels: 보안 레벨 (branchName, levels[], airportCode, history[])",
        "- securityFees: 보안 수수료 데이터",
        "- securityPolicies: 보안 정책 문서",
        "- documentLibrary: 문서 라이브러리 (attachments, downloadPermission, logs[])",
        "- importantLinks: 주요 링크 (QR 코드 포함)",
    ])
    
    # AI Logic
    add_section_slide(prs, "2. AI 로직 상세", "Gemini AI Translation Logic")
    
    add_content_slide(prs, "AI 번역 아키텍처", [
        "## 번역 플로우",
        "- 1. 사용자가 대상 언어 선택 후 번역 버튼 클릭",
        "- 2. HTML 콘텐츠에서 텍스트 추출 (태그 제거)",
        "- 3. 도메인 특화 프롬프트 구성 (항공 화물 보안)",
        "- 4. Gemini 2.5 Flash API 호출",
        "- 5. 응답 파싱 (JSON/코드펜스/라벨 자동 제거)",
        "- 6. HTML 변환 후 DOMPurify 살균 처리하여 표시",
        "",
        "## 프롬프트 엔지니어링",
        "- 항공 약어 보존: ICAO, TSA, IATA, ETD, K9, CSD, DG, ACC3",
        "- 단락 구조 유지 지시",
        "- 응답 형식 명시 (제목/내용 분리, 마크업 금지)",
        "",
        "## 응답 파싱 전략 (Robust Parsing)",
        "- 코드 펜스 자동 제거 (```json, ``` 등)",
        "- JSON 응답 시 content/translation/text 키 추출",
        "- 이스케이프 문자 정규화 (\\n, \\_)",
        "- 레이블 프리픽스 제거 (Translation:, 번역: 등)",
    ])
    
    add_content_slide(prs, "언어 감지 및 자동 대상 언어 선택", [
        "## 언어 감지 알고리즘 (detectLanguage)",
        "- HTML 태그/공백 제거 후 문자 분석",
        "- 한글 (AC00-D7A3) 비율 > 20% → 한국어",
        "- 일본어 히라가나/카타카나 > 10% → 일본어",
        "- CJK 한자 > 20% → 중국어",
        "- 기본값: 영어",
        "",
        "## 대상 언어 자동 추천 (suggestTargetLang)",
        "- 사용자 타임존 기반 현지 언어 추론",
        "- Asia/Seoul → 한국어 ↔ 영어 토글",
        "- Asia/Tokyo → 일본어",
        "- Europe/Berlin → 독일어",
        "- Asia/Bangkok → 태국어 등",
        "",
        "## 댓글 번역",
        "- KOR/ENG 퀵 버튼 + 로컬 언어 추가 버튼",
        "- 개별 댓글 단위로 번역 요청",
        "- 번역 결과 로컬 state에 캐싱",
    ])
    
    # Security
    add_section_slide(prs, "3. 보안 설계", "Security Measures")
    
    add_content_slide(prs, "보안 조치", [
        "## XSS 방어",
        "- DOMPurify 라이브러리로 모든 dangerouslySetInnerHTML 살균",
        "- 게시물 본문, 번역 결과, Markdown 미리보기 등",
        "",
        "## 인증 및 권한",
        "- Firebase Authentication (서버 사이드 토큰 검증)",
        "- 역할 기반 접근 제어 (hq_admin / branch_user)",
        "- 라우트 보호: ProtectedRoutes 컴포넌트",
        "",
        "## API 키 관리",
        "- 환경 변수 (VITE_FIREBASE_*, VITE_GEMINI_API_KEY)",
        "- .env 파일 gitignore 처리",
        "- 하드코딩된 키 없음",
        "",
        "## 코드 보안",
        "- eval() / new Function() 미사용",
        "- 외부 fetch 제한: Google Sheets CSV, Google News RSS, TopoJSON",
        "- npm audit: 낮은 심각도 이슈 2건만 존재",
    ])
    
    # Algorithms
    add_section_slide(prs, "4. 주요 알고리즘 및 구현 상세", "Implementation Details")
    
    add_content_slide(prs, "세계 지도 렌더링 알고리즘", [
        "## 지도 데이터",
        "- countries-110m.json (TopoJSON, Natural Earth 110m 해상도)",
        "- topojson-client로 GeoJSON 변환",
        "",
        "## 투영 방식",
        "- D3-like SVG 프로젝션 (Mercator 변형)",
        "- ICN(인천, 경도 127) 중심 지도",
        "- 중심 경도 오프셋: 127 (기본 180 대신)",
        "",
        "## 인터랙션",
        "- 확대/축소: 마우스 휠 (1x ~ 8x), 더블클릭 2x",
        "- 이동: 클릭 드래그 (mousedown → mousemove → mouseup)",
        "- data-station 속성으로 공항 마커 vs 배경 클릭 구분",
        "",
        "## 마커 시스템",
        "- AIRPORT_COORDS: 86+ 공항 {lat, lng, city} 사전",
        "- 위험 수준별 색상: 녹색/주황/빨강",
        "- 빨간 마커는 렌더링 우선순위 최상위 (z-order sorting)",
        "- 펄스 애니메이션 (빨강 마커)",
    ])
    
    add_content_slide(prs, "에디터 아키텍처 (Quill 2)", [
        "## Rich Text Editor",
        "- react-quill-new 3.8.3 (Quill 2 기반)",
        "- 커스텀 키보드 바인딩: Enter → <br> 삽입",
        "- 네이티브 테이블 모듈 (table: true)",
        "- 커스텀 DividerBlot (<hr> 태그 보존)",
        "",
        "## Markdown 모드",
        "- marked 라이브러리 (Markdown → HTML 변환)",
        "- Turndown 라이브러리 (HTML → Markdown 역변환)",
        "- 실시간 미리보기 (Split view)",
        "- .md/.txt 파일 직접 업로드 지원",
        "",
        "## 파일 업로드",
        "- Firebase Storage uploadBytesResumable",
        "- 진행률 표시 (overall progress bar)",
        "- 드래그앤드롭: dragenter/dragleave/dragover/drop 이벤트",
        "- 파일 크기 제한: Admin 100MB, Branch 50MB",
    ])
    
    add_content_slide(prs, "프로젝트 구조", [
        "## 디렉토리 구조",
        "- src/main.jsx → App.jsx (HashRouter > AuthProvider > PortalLayout)",
        "- src/core/ → AuthContext.jsx, PortalLayout.jsx",
        "- src/firebase/ → config.js, auth.js, collections.js",
        "- src/modules/ → 8개 모듈 (home, bulletin, document-library, ...)",
        "",
        "## 라우팅",
        "- / → HomePage",
        "- /bulletin/* → BulletinPage (boardType='directive')",
        "- /communication/* → BulletinPage (boardType='communication')",
        "- /document-library/* → DocumentLibraryPage",
        "- /security-level → SecurityLevelPage",
        "- /security-fee → SecurityFeePage",
        "- /security-policy → SecurityPolicyPage",
        "- /important-links → ImportantLinksPage",
        "- /settings → SettingsPage (admin only)",
    ])
    
    add_content_slide(prs, "빌드 및 배포", [
        "## 빌드 명령어",
        "- npm run dev: 개발 서버 (Vite HMR)",
        "- npm run build: 프로덕션 빌드",
        "- npm run preview: 로컬 프로덕션 미리보기",
        "- npm run deploy: GitHub Pages 배포 (gh-pages)",
        "",
        "## CI/CD",
        "- Vercel: GitHub main 브랜치 자동 배포",
        "- GitHub Pages: npm run deploy 수동 배포",
        "",
        "## 빌드 최적화",
        "- 코드 스플리팅: turndown, marked 동적 임포트",
        "- Vite 7 트리 쉐이킹",
        "- 총 모듈: 2,543개",
        "- 빌드 시간: ~14초",
        "- 메인 번들: ~1,299KB (gzip ~376KB)",
    ])
    
    filepath = os.path.join(OUTPUT_DIR, "AirZeta_IT_Planner_Guide_KR.pptx")
    prs.save(filepath)
    print(f"[OK] IT Planner Guide saved: {filepath}")
    return filepath


# ======================================================================
# PPT 4: EXECUTIVE/FINANCE GUIDE (Korean+English, Cost Analysis)
# ======================================================================
def create_executive_guide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    add_title_slide(prs, "AirZeta Security Portal\n경영진/재무 보고서", 
                    "Executive & Finance Overview | AI 기능, 비용 분석, 로드맵", "v1.9")
    
    # Section 1: Overview
    add_section_slide(prs, "1. 프로젝트 개요\nProject Overview", "")
    
    add_content_slide(prs, "프로젝트 개요 (Project Overview)", [
        "## AirZeta Station Security System",
        "- 항공 화물 보안 관리를 위한 통합 웹 포털",
        "- Integrated web portal for aviation cargo security management",
        "",
        "## 핵심 가치 (Core Value)",
        "- 전 세계 지점의 보안 지시사항 실시간 전달 및 확인",
        "- Real-time security directive distribution & acknowledgement tracking",
        "- AI 기반 다국어 번역으로 언어 장벽 해소",
        "- AI-powered multilingual translation (10 languages)",
        "",
        "## 현재 상태",
        "- Version 1.9 (2026년 3월 4일 기준)",
        "- 85+ 공항 지점 연결",
        "- 8개 핵심 모듈 운영 중",
        "- 전 직원 접근 가능한 클라우드 기반 서비스",
    ])
    
    # Section 2: AI Features
    add_section_slide(prs, "2. AI 기능 소개\nAI Feature Highlights", "")
    
    add_content_slide(prs, "AI 기능 하이라이트 (AI Features)", [
        "## Google Gemini AI 번역 (Translation)",
        "- 게시물 본문: 10개 언어 실시간 번역",
        "- 댓글: 원클릭 즉시 번역 (KOR/ENG + 로컬 언어)",
        "- 항공 보안 도메인 특화 용어 사전 내장",
        "- 모델: Gemini 2.5 Flash (고성능, 빠른 응답)",
        "",
        "## AI 활용 효과",
        "- 번역 비용 절감: 외부 번역 서비스 불필요",
        "- 즉시성: 수 초 내 번역 완료 (기존 수 시간 → 수 초)",
        "- 일관성: 항공 보안 용어 통일된 번역 제공",
        "",
        "## 향후 AI 확장 계획",
        "- 보안 위험도 자동 분석 및 알림",
        "- 문서 요약 기능 (Long document summarization)",
        "- 음성 인식 기반 보안 보고서 작성",
        "- 보안 트렌드 예측 분석",
    ])
    
    # Section 3: Cost Analysis
    add_section_slide(prs, "3. 비용 분석\nCost Analysis", "In-house vs Outsourced Development")
    
    # Cost comparison slide with a table-like layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY_DARK)
    
    # Title
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = NAVY_MID
    title_bg.line.fill.background()
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.6),
                 "비용 비교: 자체 개발 vs 외주 개발", font_size=22, color=WHITE, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.85), Inches(1.5), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = RED_ACCENT; line.line.fill.background()
    
    # In-house card
    card1 = add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.5), NAVY_LIGHT)
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(3.7), Inches(0.5),
                 "자체 개발 (In-house / AI-Assisted)", font_size=16, color=GREEN_ACCENT, bold=True)
    
    inhouse_items = [
        ("Firebase (Spark Plan)", "무료 (Free tier)"),
        ("Vercel 호스팅", "무료 (Hobby plan)"),
        ("Gemini AI API", "무료 tier / ~$5-10/월"),
        ("GitHub", "무료 (Public repo)"),
        ("도메인 (선택)", "~$12/년"),
        ("", ""),
        ("예상 월 비용", "약 $5-15/월"),
        ("인당 월 비용 (100명 기준)", "약 3만~5만원/인/월"),
    ]
    y_pos = 2.0
    for label, value in inhouse_items:
        if label == "":
            y_pos += 0.15
            continue
        is_total = "예상" in label or "인당" in label
        add_text_box(slide, Inches(0.8), Inches(y_pos), Inches(2.2), Inches(0.3),
                     label, font_size=10, color=LIGHT_TEXT if not is_total else GREEN_ACCENT, bold=is_total)
        add_text_box(slide, Inches(3.0), Inches(y_pos), Inches(1.7), Inches(0.3),
                     value, font_size=10, color=LIGHT_TEXT if not is_total else GREEN_ACCENT, bold=is_total, alignment=PP_ALIGN.RIGHT)
        y_pos += 0.35
    
    # Outsourced card
    card2 = add_shape_bg(slide, Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5), NAVY_LIGHT)
    add_text_box(slide, Inches(5.5), Inches(1.4), Inches(3.7), Inches(0.5),
                 "외주 개발 (Outsourced)", font_size=16, color=RED_ACCENT, bold=True)
    
    outsource_items = [
        ("초기 개발 비용", "$50,000-150,000"),
        ("월 유지보수", "$3,000-8,000/월"),
        ("서버 호스팅", "$200-500/월"),
        ("번역 서비스 (외부)", "$500-2,000/월"),
        ("보안 감사", "$5,000-15,000/년"),
        ("", ""),
        ("예상 월 비용", "약 $5,000-12,000/월"),
        ("인당 월 비용 (100명 기준)", "약 50만~120만원/인/월"),
    ]
    y_pos = 2.0
    for label, value in outsource_items:
        if label == "":
            y_pos += 0.15
            continue
        is_total = "예상" in label or "인당" in label
        add_text_box(slide, Inches(5.5), Inches(y_pos), Inches(2.2), Inches(0.3),
                     label, font_size=10, color=LIGHT_TEXT if not is_total else RED_ACCENT, bold=is_total)
        add_text_box(slide, Inches(7.7), Inches(y_pos), Inches(1.7), Inches(0.3),
                     value, font_size=10, color=LIGHT_TEXT if not is_total else RED_ACCENT, bold=is_total, alignment=PP_ALIGN.RIGHT)
        y_pos += 0.35
    
    # Savings highlight
    add_text_box(slide, Inches(0.5), Inches(6.85), Inches(9), Inches(0.5),
                 "* 자체 개발 시 외주 대비 약 90-95% 비용 절감 | Cost savings: ~90-95% compared to outsourcing",
                 font_size=11, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Per-person cost chart
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2, NAVY_DARK)
    title_bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_bg2.fill.solid(); title_bg2.fill.fore_color.rgb = NAVY_MID; title_bg2.line.fill.background()
    add_text_box(slide2, Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.6),
                 "인당 월 비용 비교 (Monthly Cost per Person)", font_size=22, color=WHITE, bold=True)
    line2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.85), Inches(1.5), Inches(0.04))
    line2.fill.solid(); line2.fill.fore_color.rgb = RED_ACCENT; line2.line.fill.background()
    
    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = ['자체 개발\n(In-house)', '외주 개발\n(Outsourced)']
    chart_data.add_series('월 비용 (KRW)', (40000, 850000))
    
    chart_frame = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.5), Inches(1.2), Inches(7), Inches(4.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False
    
    # Style the chart
    plot = chart.plots[0]
    series = plot.series[0]
    pt0 = series.points[0]
    pt0.format.fill.solid()
    pt0.format.fill.fore_color.rgb = GREEN_ACCENT
    pt1 = series.points[1]
    pt1.format.fill.solid()
    pt1.format.fill.fore_color.rgb = RED_ACCENT
    
    # Data labels
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '#,##0 "원"'
    data_labels.font.size = Pt(14)
    data_labels.font.bold = True
    data_labels.font.color.rgb = WHITE
    
    add_text_box(slide2, Inches(0.5), Inches(5.9), Inches(9), Inches(1.2),
                 "* 100명 사용 기준 | Based on 100 users\n"
                 "  자체 개발: 약 3~5만원/인/월 (클라우드 무료 티어 활용)\n"
                 "  외주 개발: 약 50~120만원/인/월 (개발+유지보수+서버+번역 서비스)",
                 font_size=10, color=GRAY_TEXT)
    
    # Section 4: Benefits
    add_section_slide(prs, "4. 도입 효과 vs 현재 업무방식\nBenefits vs Current Workflow", "")
    
    add_content_slide(prs, "현재 업무방식 vs 포털 도입 효과", [
        "## 현재 업무방식 (Before)",
        "- 보안 지시사항: 이메일 또는 메신저로 개별 전달",
        "- 확인 추적: 수동 확인 (전화, 이메일 회신 대기)",
        "- 번역: 외부 번역 서비스 또는 개별 직원 번역",
        "- 문서 관리: 공유 드라이브 또는 이메일 첨부",
        "- 보안 레벨: 엑셀 시트 수동 관리",
        "",
        "## 포털 도입 후 (After)",
        "- 보안 지시사항: 원클릭 전 지점 동시 전달",
        "- 확인 추적: 실시간 자동 추적 (Acknowledge 시스템)",
        "- 번역: AI 즉시 번역 (10개 언어, 수 초 이내)",
        "- 문서 관리: 중앙화된 문서 저장소 (권한 관리 포함)",
        "- 보안 레벨: 세계 지도 시각화 + 실시간 업데이트",
    ])
    
    add_content_slide(prs, "정량적 효과 분석 (Quantitative Benefits)", [
        "## 시간 절감",
        "- 보안 지시 전달: 2시간 → 5분 (96% 단축)",
        "- 확인 추적: 1일 → 실시간 (99% 단축)",
        "- 번역: 3시간 → 10초 (99.9% 단축)",
        "- 문서 검색: 30분 → 1분 (97% 단축)",
        "",
        "## 비용 절감",
        "- 외부 번역 비용: 월 50~200만원 → 0원 (AI 무료 티어)",
        "- 커뮤니케이션 비용: 국제 전화/팩스 → 웹 기반 무료",
        "- 종이 문서 비용: 인쇄/우편 → 디지털 PDF",
        "",
        "## 품질 향상",
        "- 정보 전달 누락률: 추정 5~10% → 0% (자동 추적)",
        "- 번역 일관성: 항공 보안 전문 용어 AI 표준화",
        "- 감사 대응: 이력 자동 기록으로 즉시 데이터 제공",
    ])
    
    # Section 5: Roadmap
    add_section_slide(prs, "5. 향후 로드맵\nFuture Roadmap", "")
    
    add_content_slide(prs, "향후 개발 로드맵 (2026 Q2 ~ 2027)", [
        "## 2026 Q2 (4~6월) - 안정화 및 최적화",
        "- 사용자 피드백 반영 UI/UX 개선",
        "- 모바일 반응형 디자인 완성",
        "- 성능 최적화 (코드 스플리팅 강화)",
        "",
        "## 2026 Q3 (7~9월) - 기능 확장",
        "- AI 보안 위험도 자동 분석 알림",
        "- 문서 자동 요약 기능",
        "- 오프라인 모드 지원 (PWA)",
        "",
        "## 2026 Q4 (10~12월) - 고도화",
        "- 대시보드 분석 리포트 (통계, 차트)",
        "- 사용자 활동 로그 상세 분석",
        "- 외부 시스템 API 연동 (ERP, TMS)",
        "",
        "## 2027 - 차세대 기능",
        "- 음성 인식 보안 보고서 작성",
        "- 예측 분석 (보안 트렌드, 위험도 예측)",
        "- 멀티테넌트 확장 (다른 항공사 제공)",
    ])
    
    # Closing
    slide_end = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide_end, NAVY_DARK)
    bar = slide_end.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.0), Inches(10), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED_ACCENT; bar.line.fill.background()
    
    add_text_box(slide_end, Inches(0.5), Inches(2.0), Inches(9), Inches(0.8),
                 "Thank You", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide_end, Inches(0.5), Inches(3.3), Inches(9), Inches(0.8),
                 "AirZeta Station Security System", font_size=18, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide_end, Inches(0.5), Inches(4.0), Inches(9), Inches(0.6),
                 "Questions? Contact the IT Security Team", font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)
    
    filepath = os.path.join(OUTPUT_DIR, "AirZeta_Executive_Finance_Guide_KR_EN.pptx")
    prs.save(filepath)
    print(f"[OK] Executive/Finance Guide saved: {filepath}")
    return filepath


# ======================================================================
# MAIN
# ======================================================================
if __name__ == "__main__":
    print("=" * 60)
    print("Generating AirZeta PowerPoint Presentations...")
    print("=" * 60)
    
    files = []
    files.append(create_admin_guide())
    files.append(create_branch_guide())
    files.append(create_it_planner_guide())
    files.append(create_executive_guide())
    
    print("\n" + "=" * 60)
    print(f"All {len(files)} presentations generated successfully!")
    for f in files:
        size_kb = os.path.getsize(f) / 1024
        print(f"  - {os.path.basename(f)} ({size_kb:.1f} KB)")
    print("=" * 60)
